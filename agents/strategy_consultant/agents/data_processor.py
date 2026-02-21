"""
agents/strategy_consultant/agents/data_processor.py
Data Processing Agent — parses Excel, Word, PPT, PDF, CSV, images.
Extracts structured content and normalizes it into context.md.
"""

import os
from pathlib import Path
from .base import BaseAgent

SYSTEM = """You are a data analyst on a consulting engagement.
You have been given extracted content from client documents.
Summarize the key data points, financials, metrics, and insights concisely.
Flag any data quality issues. Note what's missing. Structure your output clearly."""


class DataProcessorAgent(BaseAgent):
    name = "data_processor"
    default_tier = 1

    def _default_system(self): return SYSTEM

    # ── Parsers ───────────────────────────────────────────────────────────────

    def _parse_excel(self, path: Path) -> str:
        import pandas as pd
        try:
            xl = pd.ExcelFile(path)
            parts = [f"Excel file: {path.name}"]
            for sheet in xl.sheet_names[:5]:
                df = xl.parse(sheet, nrows=100)
                df = df.dropna(how="all")
                parts.append(f"\n### Sheet: {sheet} ({df.shape[0]} rows x {df.shape[1]} cols)")
                parts.append(df.to_string(max_rows=20, max_cols=10))
            return "\n".join(parts)
        except Exception as e:
            return f"[Excel parse error: {e}]"

    def _parse_csv(self, path: Path) -> str:
        import pandas as pd
        try:
            df = pd.read_csv(path, nrows=100)
            return f"CSV: {path.name} ({df.shape[0]} rows x {df.shape[1]} cols)\n{df.to_string(max_rows=20)}"
        except Exception as e:
            return f"[CSV parse error: {e}]"

    def _parse_word(self, path: Path) -> str:
        try:
            from docx import Document
            doc = Document(path)
            parts = [f"Word document: {path.name}\n"]
            for para in doc.paragraphs[:200]:
                if para.text.strip():
                    style = para.style.name
                    prefix = "## " if "Heading 1" in style else ("### " if "Heading" in style else "")
                    parts.append(f"{prefix}{para.text.strip()}")
            # Tables
            for i, table in enumerate(doc.tables[:5]):
                parts.append(f"\n[Table {i+1}]")
                for row in table.rows[:10]:
                    parts.append(" | ".join(c.text.strip() for c in row.cells))
            return "\n".join(parts)
        except Exception as e:
            return f"[Word parse error: {e}]"

    def _parse_pptx(self, path: Path) -> str:
        try:
            from pptx import Presentation
            prs = Presentation(path)
            parts = [f"PowerPoint: {path.name} ({len(prs.slides)} slides)\n"]
            for i, slide in enumerate(prs.slides[:30], 1):
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text.strip())
                if texts:
                    parts.append(f"**Slide {i}:** " + " | ".join(texts[:5]))
                # Speaker notes
                if slide.has_notes_slide:
                    note = slide.notes_slide.notes_text_frame.text.strip()
                    if note:
                        parts.append(f"  [Notes]: {note[:200]}")
            return "\n".join(parts)
        except Exception as e:
            return f"[PPTX parse error: {e}]"

    def _parse_pdf(self, path: Path) -> str:
        try:
            import pdfplumber
            parts = [f"PDF: {path.name}\n"]
            with pdfplumber.open(path) as pdf:
                for i, page in enumerate(pdf.pages[:20], 1):
                    text = page.extract_text() or ""
                    if text.strip():
                        parts.append(f"\n[Page {i}]\n{text[:1500]}")
                    tables = page.extract_tables()
                    for t in tables[:2]:
                        parts.append(f"\n[Table on page {i}]")
                        for row in t[:8]:
                            parts.append(" | ".join(str(c or "") for c in row))
            return "\n".join(parts)
        except Exception as e:
            return f"[PDF parse error: {e}]"

    def _parse_image(self, path: Path) -> str:
        # Use vision model via litellm
        try:
            import base64, litellm
            with open(path, "rb") as f:
                b64 = base64.b64encode(f.read()).decode()
            ext = path.suffix.lower().lstrip(".")
            mime = {"jpg": "jpeg", "jpeg": "jpeg", "png": "png", "gif": "gif", "webp": "webp"}.get(ext, "jpeg")
            resp = litellm.completion(
                model="anthropic/claude-haiku-4-5",
                messages=[{"role": "user", "content": [
                    {"type": "image_url", "image_url": {"url": f"data:image/{mime};base64,{b64}"}},
                    {"type": "text", "text": "Describe this image in detail. Extract any text, numbers, charts, or diagrams visible."}
                ]}],
                max_tokens=500,
            )
            return f"Image: {path.name}\n{resp.choices[0].message.content}"
        except Exception as e:
            return f"[Image parse error: {e}]"

    def parse_file(self, path: Path) -> str:
        """Dispatch to correct parser by extension."""
        ext = path.suffix.lower()
        dispatch = {
            ".xlsx": self._parse_excel, ".xls": self._parse_excel,
            ".csv":  self._parse_csv,
            ".docx": self._parse_word, ".doc": self._parse_word,
            ".pptx": self._parse_pptx, ".ppt": self._parse_pptx,
            ".pdf":  self._parse_pdf,
            ".png":  self._parse_image, ".jpg": self._parse_image,
            ".jpeg": self._parse_image, ".webp": self._parse_image,
        }
        parser = dispatch.get(ext)
        if parser:
            raw = parser(path)
            # Summarize with LLM if long
            if len(raw) > 3000:
                prompt = f"Summarize the key data, metrics, and insights from this document content:\n\n{raw[:6000]}"
                return self.call(prompt, force_tier=1)
            return raw
        return f"[Unsupported format: {ext}]"

    def run(self, data_dir: Path = None, files: list = None) -> dict:
        """
        Process all files in data_dir or the given file list.
        Returns {filename: extracted_content} dict.
        """
        results = {}
        targets = []

        if files:
            targets = [Path(f) for f in files]
        elif data_dir:
            targets = [f for f in Path(data_dir).iterdir() if f.is_file()]

        for path in targets:
            self.logger.info(f"Processing: {path.name}")
            content = self.parse_file(path)
            results[path.name] = content

        return results

    def validate(self) -> dict:
        """Test with a synthetic CSV in memory."""
        import tempfile, csv
        try:
            with tempfile.NamedTemporaryFile(mode="w", suffix=".csv", delete=False, newline="") as f:
                writer = csv.writer(f)
                writer.writerow(["Company", "Revenue_M", "Growth_pct", "Employees"])
                writer.writerows([
                    ["AlphaAI", 120, 45, 850],
                    ["BetaML", 80, 62, 420],
                    ["GammaData", 200, 28, 1200],
                ])
                tmp = f.name
            result = self.run(files=[tmp])
            os.unlink(tmp)
            passed = len(result) == 1 and "AlphaAI" in list(result.values())[0]
            return {
                "passed": passed,
                "output": list(result.values())[0][:300] if result else "",
                "notes": "CSV parsing test with synthetic competitive data"
            }
        except Exception as e:
            return {"passed": False, "output": "", "notes": str(e)}
